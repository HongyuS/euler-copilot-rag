"""
PowerPoint 文件解析器
"""
import logging
from typing import Optional
from pptx import Presentation
from pptx.table import Table

logger = logging.getLogger(__name__)


class PptxParser:
    """PowerPoint 文档解析器"""
    
    @staticmethod
    def _extract_table_to_array(table: Table) -> list:
        """提取表格为数组"""
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = ''.join([p.text for p in cell.text_frame.paragraphs])
                row_data.append(cell_text)
            table_data.append(row_data)
        return table_data
    
    @staticmethod
    async def parse(file_path: str) -> Optional[str]:
        """
        解析 PowerPoint 文件
        
        :param file_path: 文件路径
        :return: 文件内容
        """
        try:
            pptx = Presentation(file_path)
            if not pptx:
                logger.error("[PptxParser] 无法打开 PPTX 文件")
                return None
            
            paragraphs = []
            
            for slide_num, slide in enumerate(pptx.slides, start=1):
                paragraphs.append(f"\n[幻灯片 {slide_num}]")
                
                for shape in slide.shapes:
                    # 提取文字
                    if shape.has_text_frame:
                        text = ""
                        try:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    text += run.text
                        except Exception as e:
                            logger.warning(f"[PptxParser] 文字提取失败: {e}")
                        
                        if text.strip():
                            paragraphs.append(text)
                    
                    # 提取表格
                    elif shape.has_table:
                        table = shape.table
                        table_array = PptxParser._extract_table_to_array(table)
                        for row in table_array:
                            paragraphs.append(' | '.join(row))
                    
                    # 图片标记
                    elif shape.shape_type == 13:  # 13 表示图片类型
                        paragraphs.append("[图片]")
            
            content = '\n'.join(paragraphs)
            return content if content.strip() else None
        except Exception as e:
            logger.exception(f"[PptxParser] 解析 PPTX 文件失败: {e}")
            return None

